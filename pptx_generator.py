"""
pptx_generator.py
Gera a apresentação PPTX de 7 slides com base nos dados do Claude.
Design: Midnight Executive — Navy / Ice Blue / White
"""

import io
import logging
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import lxml.etree as etree

logger = logging.getLogger(__name__)

# ── Paleta ───────────────────────────────────────────────────────────────────
C_DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
C_MID_NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
C_ACCENT_BLUE = RGBColor(0x1F, 0x5E, 0xFF)
C_ICE_BLUE    = RGBColor(0xCA, 0xDC, 0xFC)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xFA)
C_MID_GREY    = RGBColor(0x8A, 0x97, 0xA8)
C_RED         = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN       = RGBColor(0x1A, 0x7A, 0x4A)
C_AMBER       = RGBColor(0xD4, 0xA0, 0x17)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _rgb_hex(r: RGBColor) -> str:
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, alpha=None):
    """Adiciona retângulo colorido."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def _add_text_box(slide, text, left, top, width, height,
                  font_name="Calibri", font_size=14, bold=False, italic=False,
                  color: RGBColor = C_WHITE, align=PP_ALIGN.LEFT,
                  word_wrap=True):
    """Adiciona text box formatado."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _add_multi_line_box(slide, lines, left, top, width, height,
                        font_name="Calibri", font_size=11,
                        color: RGBColor = C_WHITE,
                        bold_first=False,
                        spacing_after=Pt(4)):
    """Adiciona caixa de texto com múltiplas linhas."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing_after
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
    return txBox


def _slide_number_label(slide, number: int):
    """Adiciona rótulo de número de slide no canto inferior direito."""
    _add_text_box(
        slide, str(number),
        SLIDE_W - Inches(0.5), SLIDE_H - Inches(0.35),
        Inches(0.4), Inches(0.3),
        font_size=9, color=C_MID_GREY, align=PP_ALIGN.RIGHT
    )


def _footer_bar(slide, label: str):
    """Rodapé discreto."""
    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), C_MID_NAVY)
    _add_text_box(
        slide, label,
        Inches(0.3), SLIDE_H - Inches(0.28),
        SLIDE_W - Inches(0.6), Inches(0.26),
        font_size=7.5, color=C_MID_GREY
    )


def _fmt(val, prefix="", suffix="", decimals=2, na="N/A"):
    if val is None:
        return na
    try:
        return f"{prefix}{float(val):.{decimals}f}{suffix}"
    except (TypeError, ValueError):
        return str(val) if val else na


def _rating_color(rating: str) -> RGBColor:
    r = (rating or "").upper()
    if r == "BUY":
        return C_GREEN
    elif r == "SELL":
        return C_RED
    return C_AMBER


def _add_table(slide, rows, col_widths, left, top, row_height=Inches(0.38),
               header_fill=C_MID_NAVY, alt_fill=C_LIGHT_GREY,
               font_size=10):
    """Cria uma tabela formatada no slide."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths)
    total_h = row_height * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, total_h).table
    table.first_row = True

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)

            # Background
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{_rgb_hex(header_fill if ri == 0 else (C_LIGHT_GREY if ri % 2 == 0 else C_WHITE))}"/>'
                f'</a:solidFill>'
            )
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            tcPr.append(solidFill)

            # Text formatting
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            if para.runs:
                run = para.runs[0]
            else:
                run = para.add_run()
                run.text = str(cell_text)
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.bold = (ri == 0)
            run.font.color.rgb = C_WHITE if ri == 0 else C_DARK_NAVY

    return table


# ═════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════════════════

def _slide_1_cover(prs, analysis, ticker):
    """Slide 1: Capa."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    company = analysis.get("company_summary", {})
    verdict = analysis.get("verdict", {})
    rating = verdict.get("rating", "N/A")
    rc = _rating_color(rating)

    # Full dark background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_DARK_NAVY)

    # Left accent stripe
    _add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, C_ACCENT_BLUE)

    # Company name
    _add_text_box(
        slide, company.get("name", ticker),
        Inches(0.5), Inches(1.2),
        Inches(8), Inches(1.4),
        font_name="Calibri", font_size=40, bold=True,
        color=C_WHITE, align=PP_ALIGN.LEFT
    )

    # Ticker & sector
    sub = f"{ticker}  |  {company.get('sector', '')}  |  {company.get('industry', '')}"
    _add_text_box(
        slide, sub,
        Inches(0.5), Inches(2.6),
        Inches(9), Inches(0.5),
        font_size=14, color=C_ICE_BLUE
    )

    # Separator line
    _add_rect(slide, Inches(0.5), Inches(3.25), Inches(8), Inches(0.025), C_ACCENT_BLUE)

    # Subtitle
    _add_text_box(
        slide, "INSTITUTIONAL EQUITY RESEARCH",
        Inches(0.5), Inches(3.45),
        Inches(8), Inches(0.5),
        font_size=12, color=C_MID_GREY, bold=False
    )

    # Rating badge
    _add_rect(slide, Inches(9.5), Inches(1.2), Inches(3.0), Inches(1.1), rc)
    _add_text_box(
        slide, "RATING",
        Inches(9.5), Inches(1.25),
        Inches(3.0), Inches(0.35),
        font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER
    )
    _add_text_box(
        slide, rating,
        Inches(9.5), Inches(1.55),
        Inches(3.0), Inches(0.65),
        font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER
    )

    # Target price box
    tp = _fmt(verdict.get("blended_target_price"), prefix="$")
    up = _fmt(verdict.get("upside_pct"), suffix="%")
    cp = _fmt(verdict.get("current_price"), prefix="$")

    _add_rect(slide, Inches(9.5), Inches(2.5), Inches(3.0), Inches(2.0), C_MID_NAVY)
    _add_text_box(slide, "Target Price", Inches(9.5), Inches(2.6), Inches(3.0), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, tp, Inches(9.5), Inches(2.95), Inches(3.0), Inches(0.65),
                  font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, f"Current: {cp}  |  Upside: {up}",
                  Inches(9.5), Inches(3.6), Inches(3.0), Inches(0.4),
                  font_size=9, color=C_MID_GREY, align=PP_ALIGN.CENTER)

    # Date
    _add_text_box(
        slide, datetime.utcnow().strftime("%B %d, %Y"),
        Inches(0.5), Inches(6.8),
        Inches(5), Inches(0.35),
        font_size=9, color=C_MID_GREY
    )

    _footer_bar(slide, "CONFIDENTIAL — For Institutional Use Only")
    _slide_number_label(slide, 1)


def _slide_2_business(prs, analysis):
    """Slide 2: Business Model & Tese."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    biz = analysis.get("company_summary", {})

    # Header bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_MID_NAVY)
    _add_text_box(slide, "Business Model & Investment Thesis",
                  Inches(0.3), Inches(0.2), Inches(9), Inches(0.7),
                  font_size=24, bold=True, color=C_WHITE)
    _add_text_box(slide, "01 / 07",
                  Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.RIGHT)

    # Left column: Business Model
    _add_rect(slide, Inches(0.3), Inches(1.25), Inches(5.8), Inches(0.38), C_ACCENT_BLUE)
    _add_text_box(slide, "BUSINESS MODEL",
                  Inches(0.35), Inches(1.27), Inches(5.7), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)

    _add_text_box(slide, biz.get("business_model", "N/A"),
                  Inches(0.3), Inches(1.7), Inches(5.8), Inches(2.0),
                  font_size=10.5, color=C_DARK_NAVY, word_wrap=True)

    # Revenue streams
    _add_rect(slide, Inches(0.3), Inches(3.85), Inches(5.8), Inches(0.38), C_MID_NAVY)
    _add_text_box(slide, "REVENUE STREAMS",
                  Inches(0.35), Inches(3.87), Inches(5.7), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)

    rev_lines = [f"• {rs}" for rs in biz.get("revenue_streams", ["N/A"])[:4]]
    _add_multi_line_box(slide, rev_lines, Inches(0.3), Inches(4.3),
                        Inches(5.8), Inches(1.8), font_size=10.5, color=C_DARK_NAVY)

    # Right column: Moat
    _add_rect(slide, Inches(6.7), Inches(1.25), Inches(6.0), Inches(0.38), C_MID_NAVY)
    _add_text_box(slide, "COMPETITIVE MOAT",
                  Inches(6.75), Inches(1.27), Inches(5.9), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)

    moat_lines = [f"→ {mf}" for mf in biz.get("moat_factors", ["N/A"])[:5]]
    _add_multi_line_box(slide, moat_lines, Inches(6.7), Inches(1.7),
                        Inches(6.0), Inches(2.1), font_size=10.5, color=C_DARK_NAVY)

    # Moat assessment
    _add_rect(slide, Inches(6.7), Inches(3.85), Inches(6.0), Inches(0.38), C_DARK_NAVY)
    _add_text_box(slide, "MOAT DURABILITY",
                  Inches(6.75), Inches(3.87), Inches(5.9), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)
    _add_text_box(slide, biz.get("moat_assessment", "N/A"),
                  Inches(6.7), Inches(4.3), Inches(6.0), Inches(1.9),
                  font_size=10.5, color=C_DARK_NAVY, word_wrap=True)

    _footer_bar(slide, "CONFIDENTIAL — Institutional Research")
    _slide_number_label(slide, 2)


def _slide_3_technical(prs, analysis):
    """Slide 3: Análise Técnica & Insiders."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tech = analysis.get("technical_analysis", {})

    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_MID_NAVY)
    _add_text_box(slide, "Technical Analysis & Insider Activity",
                  Inches(0.3), Inches(0.2), Inches(9), Inches(0.7),
                  font_size=24, bold=True, color=C_WHITE)
    _add_text_box(slide, "02 / 07",
                  Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.RIGHT)

    # KPI boxes top row
    kpis = [
        ("Current Price", _fmt(tech.get("current_price"), prefix="$")),
        ("MA 50", _fmt(tech.get("ma50"), prefix="$")),
        ("MA 200", _fmt(tech.get("ma200"), prefix="$")),
        ("RSI (14)", _fmt(tech.get("rsi_14"))),
    ]
    box_w = Inches(2.9)
    for i, (label, val) in enumerate(kpis):
        x = Inches(0.3) + i * (box_w + Inches(0.25))
        _add_rect(slide, x, Inches(1.3), box_w, Inches(1.1), C_LIGHT_GREY)
        _add_text_box(slide, label, x, Inches(1.35), box_w, Inches(0.38),
                      font_size=9.5, color=C_MID_GREY, align=PP_ALIGN.CENTER)
        _add_text_box(slide, val, x, Inches(1.72), box_w, Inches(0.55),
                      font_size=20, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

    # Signal table
    sig_data = [
        ["Indicator", "Signal / Reading"],
        ["RSI Signal", tech.get("rsi_signal", "N/A")],
        ["MACD Signal", tech.get("macd_signal", "N/A")],
        ["Trend Bias", tech.get("trend_bias", "N/A")],
        ["Price vs MA50", _fmt(tech.get("price_vs_ma50_pct"), suffix="%")],
        ["Price vs MA200", _fmt(tech.get("price_vs_ma200_pct"), suffix="%")],
    ]
    _add_table(slide, sig_data,
               [Inches(2.8), Inches(3.5)],
               left=Inches(0.3), top=Inches(2.7),
               row_height=Inches(0.37), font_size=10)

    # Insider section
    _add_rect(slide, Inches(7.0), Inches(2.7), Inches(5.9), Inches(0.38), C_DARK_NAVY)
    _add_text_box(slide, "INSIDER ACTIVITY",
                  Inches(7.05), Inches(2.72), Inches(5.8), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)
    _add_text_box(slide, f"Pattern: {tech.get('insider_pattern', 'N/A')}",
                  Inches(7.0), Inches(3.15), Inches(5.9), Inches(0.45),
                  font_size=12, bold=True, color=C_MID_NAVY)
    _add_text_box(slide, tech.get("insider_summary", "N/A"),
                  Inches(7.0), Inches(3.6), Inches(5.9), Inches(2.5),
                  font_size=10.5, color=C_DARK_NAVY, word_wrap=True)

    _footer_bar(slide, "CONFIDENTIAL — Institutional Research")
    _slide_number_label(slide, 3)


def _slide_4_dcf(prs, analysis):
    """Slide 4: Modelo DCF."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dcf = analysis.get("dcf_model", {})

    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_MID_NAVY)
    _add_text_box(slide, "DCF Valuation Model — Mid-Year Convention",
                  Inches(0.3), Inches(0.2), Inches(10), Inches(0.7),
                  font_size=24, bold=True, color=C_WHITE)
    _add_text_box(slide, "03 / 07",
                  Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.RIGHT)

    # Assumptions box
    _add_rect(slide, Inches(0.3), Inches(1.25), Inches(4.2), Inches(5.8), C_LIGHT_GREY)
    _add_text_box(slide, "MODEL ASSUMPTIONS",
                  Inches(0.35), Inches(1.3), Inches(4.1), Inches(0.4),
                  font_size=10, bold=True, color=C_MID_NAVY)

    g = dcf.get("growth_assumptions", {})
    assump_lines = [
        f"Base FCF: {_fmt(dcf.get('base_fcf_bn'), suffix='B')}",
        f"WACC: {_fmt(dcf.get('wacc_pct'), suffix='%')}",
        f"Terminal Growth: {_fmt(dcf.get('terminal_growth_rate_pct'), suffix='%')}",
        "",
        f"Growth — Conservative: {_fmt(g.get('conservative'), suffix='%')}",
        f"Growth — Base: {_fmt(g.get('base'), suffix='%')}",
        f"Growth — Bull: {_fmt(g.get('bull'), suffix='%')}",
        "",
        f"Net Debt: {_fmt(dcf.get('net_debt_bn'), suffix='B')}",
        f"Shares Out: {_fmt(dcf.get('shares_outstanding_bn'), decimals=3)}B",
    ]
    _add_multi_line_box(slide, assump_lines, Inches(0.4), Inches(1.75),
                        Inches(3.9), Inches(4.8), font_size=11, color=C_DARK_NAVY)

    # FCF projections table
    proj = dcf.get("fcf_projections_bn", [])
    pvs  = dcf.get("pv_fcfs_bn", [])
    if proj:
        n = len(proj)
        yr_headers = [""] + [f"Y{i+1}" for i in range(n)]
        fcf_row = ["FCF ($B)"] + [_fmt(v) for v in proj]
        pv_row  = ["PV ($B)"]  + [_fmt(v) for v in pvs]
        proj_data = [yr_headers, fcf_row, pv_row]
        col_w_proj = [Inches(1.2)] + [Inches(1.45)] * n
        if n <= 5:
            _add_table(slide, proj_data, col_w_proj,
                       left=Inches(4.8), top=Inches(1.3),
                       row_height=Inches(0.42), font_size=10)

    # Valuation bridge
    bridge_data = [
        ["Component", "Value ($B)"],
        ["Sum PV FCFs", _fmt(dcf.get("sum_pv_fcfs_bn"), suffix="B")],
        ["PV Terminal Value", _fmt(dcf.get("pv_terminal_value_bn"), suffix="B")],
        ["Enterprise Value", _fmt(dcf.get("enterprise_value_bn"), suffix="B")],
        ["Less Net Debt", _fmt(dcf.get("net_debt_bn"), suffix="B")],
        ["Equity Value", _fmt(dcf.get("equity_value_bn"), suffix="B")],
    ]
    _add_table(slide, bridge_data,
               [Inches(3.5), Inches(2.3)],
               left=Inches(4.8), top=Inches(3.0),
               row_height=Inches(0.4), font_size=10)

    # Intrinsic value highlight
    iv = _fmt(dcf.get("dcf_intrinsic_value"), prefix="$")
    _add_rect(slide, Inches(4.8), Inches(5.7), Inches(5.8), Inches(0.9), C_ACCENT_BLUE)
    _add_text_box(slide, "DCF INTRINSIC VALUE PER SHARE",
                  Inches(4.85), Inches(5.72), Inches(5.7), Inches(0.35),
                  font_size=9, color=C_ICE_BLUE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, iv, Inches(4.85), Inches(6.05), Inches(5.7), Inches(0.5),
                  font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _footer_bar(slide, "CONFIDENTIAL — Institutional Research")
    _slide_number_label(slide, 4)


def _slide_5_multiples(prs, analysis, ticker):
    """Slide 5: Múltiplos Comparativos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mult = analysis.get("multiples_analysis", {})
    subj = mult.get("subject", {})
    p1   = mult.get("peer_1", {})
    p2   = mult.get("peer_2", {})

    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_MID_NAVY)
    _add_text_box(slide, "Comparable Company Multiples Analysis",
                  Inches(0.3), Inches(0.2), Inches(10), Inches(0.7),
                  font_size=24, bold=True, color=C_WHITE)
    _add_text_box(slide, "04 / 07",
                  Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.RIGHT)

    comp_data = [
        ["Metric", ticker, f"{p1.get('name', 'Peer 1')}", f"{p2.get('name', 'Peer 2')}"],
        ["EV/EBITDA",    _fmt(subj.get("ev_ebitda")),    _fmt(p1.get("ev_ebitda")),  _fmt(p2.get("ev_ebitda"))],
        ["P/E (TTM)",    _fmt(subj.get("pe_ttm")),       _fmt(p1.get("pe_ttm")),     _fmt(p2.get("pe_ttm"))],
        ["P/S (TTM)",    _fmt(subj.get("ps_ttm")),       _fmt(p1.get("ps_ttm")),     _fmt(p2.get("ps_ttm"))],
        ["Rev Growth %", _fmt(subj.get("revenue_growth_yoy_pct"), suffix="%"), "—", "—"],
        ["Value/Growth", _fmt(subj.get("value_growth_score"), decimals=2), "—", "—"],
    ]
    _add_table(slide, comp_data,
               [Inches(2.5), Inches(2.8), Inches(2.8), Inches(2.8)],
               left=Inches(1.0), top=Inches(1.4),
               row_height=Inches(0.5), font_size=11)

    # Multiples implied price
    _add_rect(slide, Inches(1.0), Inches(5.3), Inches(5.5), Inches(0.9), C_MID_NAVY)
    _add_text_box(slide, "Multiples-Implied Target Price",
                  Inches(1.05), Inches(5.32), Inches(5.4), Inches(0.38),
                  font_size=9.5, color=C_ICE_BLUE, align=PP_ALIGN.CENTER)
    _add_text_box(slide,
                  _fmt(mult.get("multiples_implied_price"), prefix="$"),
                  Inches(1.05), Inches(5.65), Inches(5.4), Inches(0.5),
                  font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Methodology
    meth = mult.get("multiples_methodology", "")
    if meth:
        _add_text_box(slide, meth,
                      Inches(7.2), Inches(1.4), Inches(5.7), Inches(4.5),
                      font_size=10.5, color=C_DARK_NAVY, word_wrap=True)

    _footer_bar(slide, "CONFIDENTIAL — Institutional Research")
    _slide_number_label(slide, 5)


def _slide_6_risks(prs, analysis):
    """Slide 6: Riscos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bear = analysis.get("bear_case", {})

    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_MID_NAVY)
    _add_text_box(slide, "Bear Case — Principal Risk Factors",
                  Inches(0.3), Inches(0.2), Inches(10), Inches(0.7),
                  font_size=24, bold=True, color=C_WHITE)
    _add_text_box(slide, "05 / 07",
                  Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.RIGHT)

    risk_keys = ["risk_1", "risk_2", "risk_3"]
    y_positions = [Inches(1.3), Inches(3.1), Inches(4.9)]

    for i, (rk, y_pos) in enumerate(zip(risk_keys, y_positions)):
        risk = bear.get(rk, {})
        if not risk:
            continue

        _add_rect(slide, Inches(0.3), y_pos, Inches(12.7), Inches(1.65), C_LIGHT_GREY)

        # Category badge
        _add_rect(slide, Inches(0.3), y_pos, Inches(0.12), Inches(1.65), C_RED)
        _add_text_box(slide, f"{i+1}. {risk.get('category', 'Risk')}",
                      Inches(0.55), y_pos + Inches(0.08), Inches(5), Inches(0.42),
                      font_size=11, bold=True, color=C_RED)

        # Probability / Impact chips
        prob_col = C_AMBER if risk.get("probability", "").lower() == "medium" else (
            C_RED if risk.get("probability", "").lower() == "high" else C_GREEN
        )
        _add_rect(slide, Inches(7.8), y_pos + Inches(0.08), Inches(1.5), Inches(0.32), prob_col)
        _add_text_box(slide, f"Prob: {risk.get('probability', 'N/A')}",
                      Inches(7.82), y_pos + Inches(0.09), Inches(1.46), Inches(0.28),
                      font_size=8.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        imp_col = C_RED if risk.get("impact", "").lower() == "high" else (
            C_AMBER if risk.get("impact", "").lower() == "medium" else C_GREEN
        )
        _add_rect(slide, Inches(9.5), y_pos + Inches(0.08), Inches(1.5), Inches(0.32), imp_col)
        _add_text_box(slide, f"Impact: {risk.get('impact', 'N/A')}",
                      Inches(9.52), y_pos + Inches(0.09), Inches(1.46), Inches(0.28),
                      font_size=8.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Description
        _add_text_box(slide, risk.get("description", "N/A"),
                      Inches(0.55), y_pos + Inches(0.52), Inches(12.1), Inches(0.6),
                      font_size=10.5, color=C_DARK_NAVY, word_wrap=True)

        # Mitigant
        _add_text_box(slide, f"Mitigant: {risk.get('mitigant', 'N/A')}",
                      Inches(0.55), y_pos + Inches(1.18), Inches(12.1), Inches(0.38),
                      font_size=9.5, color=C_MID_NAVY, word_wrap=True)

    _footer_bar(slide, "CONFIDENTIAL — Institutional Research")
    _slide_number_label(slide, 6)


def _slide_7_verdict(prs, analysis, ticker):
    """Slide 7: Conclusão e Preço-Alvo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    verdict = analysis.get("verdict", {})
    rating = verdict.get("rating", "N/A")
    rc = _rating_color(rating)

    # Full dark background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_DARK_NAVY)
    _add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, rc)

    _add_text_box(slide, "Verdict & Price Target",
                  Inches(0.4), Inches(0.3), Inches(9), Inches(0.8),
                  font_size=28, bold=True, color=C_WHITE)
    _add_text_box(slide, "06 / 07",
                  Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.4),
                  font_size=10, color=C_MID_GREY, align=PP_ALIGN.RIGHT)

    # Three price boxes
    boxes = [
        ("DCF Intrinsic Value", _fmt(verdict.get("dcf_target_price"), prefix="$"), "50% Weight"),
        ("Multiples Target", _fmt(verdict.get("multiples_target_price"), prefix="$"), "50% Weight"),
        ("Blended Target", _fmt(verdict.get("blended_target_price"), prefix="$"), "Final Price Target"),
    ]
    box_w = Inches(3.5)
    for i, (label, price, sub) in enumerate(boxes):
        x = Inches(0.4) + i * (box_w + Inches(0.35))
        fill = C_ACCENT_BLUE if i == 2 else C_MID_NAVY
        _add_rect(slide, x, Inches(1.4), box_w, Inches(1.6), fill)
        _add_text_box(slide, label, x, Inches(1.45), box_w, Inches(0.45),
                      font_size=10, color=C_ICE_BLUE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, price, x, Inches(1.85), box_w, Inches(0.7),
                      font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, sub, x, Inches(2.55), box_w, Inches(0.35),
                      font_size=9, color=C_MID_GREY, align=PP_ALIGN.CENTER)

    # Upside & Rating
    cp = _fmt(verdict.get("current_price"), prefix="$")
    up = _fmt(verdict.get("upside_pct"), suffix="%")
    _add_text_box(slide, f"Current Price: {cp}    Upside / Downside: {up}",
                  Inches(0.4), Inches(3.25), Inches(8), Inches(0.45),
                  font_size=12, color=C_ICE_BLUE)

    # Rating badge
    _add_rect(slide, Inches(10.0), Inches(1.4), Inches(3.0), Inches(1.6), rc)
    _add_text_box(slide, "FINAL RATING", Inches(10.0), Inches(1.5), Inches(3.0), Inches(0.45),
                  font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_text_box(slide, rating, Inches(10.0), Inches(1.9), Inches(3.0), Inches(0.9),
                  font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Investment thesis
    _add_rect(slide, Inches(0.4), Inches(3.85), Inches(12.5), Inches(0.38), C_MID_NAVY)
    _add_text_box(slide, "INVESTMENT THESIS",
                  Inches(0.45), Inches(3.87), Inches(12.3), Inches(0.34),
                  font_size=10, bold=True, color=C_WHITE)
    thesis = verdict.get("investment_thesis", "N/A")
    _add_text_box(slide, thesis,
                  Inches(0.4), Inches(4.3), Inches(12.5), Inches(1.9),
                  font_size=11, color=C_ICE_BLUE, word_wrap=True)

    # Disclaimer
    _add_text_box(slide,
                  "This analysis is for institutional use only and does not constitute investment advice.",
                  Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
                  font_size=7.5, color=C_MID_GREY)

    _slide_number_label(slide, 7)


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def generate_pptx(analysis: dict, ticker: str) -> bytes:
    """
    Gera a apresentação PPTX de 7 slides e retorna os bytes.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_1_cover(prs, analysis, ticker)
    _slide_2_business(prs, analysis)
    _slide_3_technical(prs, analysis)
    _slide_4_dcf(prs, analysis)
    _slide_5_multiples(prs, analysis, ticker)
    _slide_6_risks(prs, analysis)
    _slide_7_verdict(prs, analysis, ticker)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
